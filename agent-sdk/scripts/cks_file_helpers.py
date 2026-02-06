# -*- coding: utf-8 -*-
"""
CKS Lite - File Operation Helpers
Reusable utilities for document manipulation tasks.

The AI can copy this file to the user's temp dir and import it,
or use the functions as reference for inline scripts.
"""

import os
import sys

# Ensure UTF-8 output on Windows
if sys.platform == "win32":
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")


def get_temp_dir():
    """Get a temp directory for CKS scripts."""
    temp = os.path.join(os.environ.get("TEMP", os.environ.get("TMP", "/tmp")), "cks_lite")
    os.makedirs(temp, exist_ok=True)
    return temp


def read_pdf(path):
    """Read PDF and return text content per page."""
    import fitz
    doc = fitz.open(path)
    pages = []
    for i in range(len(doc)):
        pages.append({"page": i + 1, "text": doc[i].get_text()})
    doc.close()
    return pages


def read_excel(path, max_rows=100):
    """Read Excel and return sheet data."""
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    sheets = {}
    for name in wb.sheetnames:
        ws = wb[name]
        rows = []
        for row in ws.iter_rows(max_row=min(ws.max_row, max_rows), values_only=True):
            rows.append([str(c) if c is not None else "" for c in row])
        sheets[name] = {"rows": ws.max_row, "cols": ws.max_column, "data": rows}
    return sheets


def read_docx(path):
    """Read Word document and return structured content."""
    from docx import Document
    doc = Document(path)
    content = []
    for para in doc.paragraphs:
        if para.text.strip():
            content.append({"style": para.style.name, "text": para.text})
    tables = []
    for table in doc.tables:
        table_data = []
        for row in table.rows:
            table_data.append([cell.text for cell in row.cells])
        tables.append(table_data)
    return {"paragraphs": content, "tables": tables}


def read_pptx(path):
    """Read PowerPoint and return slide content."""
    from pptx import Presentation
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
        slides.append({"slide": i + 1, "texts": texts})
    return slides


def create_chart(data, labels, title, output_path, chart_type="bar"):
    """Create a chart image from data."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False

    fig, ax = plt.subplots(figsize=(10, 6))

    if chart_type == "bar":
        ax.bar(labels, data)
    elif chart_type == "pie":
        ax.pie(data, labels=labels, autopct="%1.1f%%")
    elif chart_type == "line":
        ax.plot(labels, data, marker="o")

    ax.set_title(title)
    fig.savefig(output_path, dpi=150, bbox_inches="tight")
    plt.close(fig)
    return output_path


if __name__ == "__main__":
    print("CKS File Helpers loaded. Available functions:")
    print("  read_pdf(path)")
    print("  read_excel(path, max_rows=100)")
    print("  read_docx(path)")
    print("  read_pptx(path)")
    print("  create_chart(data, labels, title, output_path)")
    print(f"  Temp dir: {get_temp_dir()}")
