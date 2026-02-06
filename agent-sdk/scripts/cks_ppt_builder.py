# -*- coding: utf-8 -*-
"""
CKS Lite PPT Builder - 预构建 PPT 生成器
用法: python cks_ppt_builder.py <config.json>

config.json 格式:
{
    "title": "演示标题",
    "subtitle": "副标题 | 日期",
    "style": "business",  // business | tech | minimal | creative
    "slides": [
        {
            "title": "章节标题",
            "bullets": ["要点1（详细描述）", "要点2", ...],
            "notes": "演讲者备注（可选）"
        },
        ...
    ],
    "output": "C:\\path\\to\\output.pptx"
}
"""

import json
import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx 未安装，正在安装...")
    os.system(f"{sys.executable} -m pip install python-pptx -q")
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor


# ============ 颜色方案 ============

STYLES = {
    "business": {
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0x00, 0xD2, 0xFF),
        "accent2": RGBColor(0x7C, 0x3A, 0xED),
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtext": RGBColor(0xB0, 0xB0, 0xB0),
        "divider": RGBColor(0x33, 0x33, 0x55),
    },
    "tech": {
        "bg": RGBColor(0x0D, 0x0D, 0x0D),
        "accent": RGBColor(0x00, 0xFF, 0x88),
        "accent2": RGBColor(0x00, 0xBB, 0xFF),
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtext": RGBColor(0xA0, 0xA0, 0xA0),
        "divider": RGBColor(0x2A, 0x2A, 0x2A),
    },
    "minimal": {
        "bg": RGBColor(0xFA, 0xFA, 0xFA),
        "accent": RGBColor(0x1A, 0x1A, 0x1A),
        "accent2": RGBColor(0x66, 0x66, 0x66),
        "text": RGBColor(0x1A, 0x1A, 0x1A),
        "subtext": RGBColor(0x88, 0x88, 0x88),
        "divider": RGBColor(0xDD, 0xDD, 0xDD),
    },
    "creative": {
        "bg": RGBColor(0x1A, 0x0A, 0x2E),
        "accent": RGBColor(0xFF, 0x6B, 0x9D),
        "accent2": RGBColor(0x7C, 0x3A, 0xED),
        "text": RGBColor(0xFF, 0xFF, 0xFF),
        "subtext": RGBColor(0xC0, 0xA0, 0xD0),
        "divider": RGBColor(0x3A, 0x2A, 0x55),
    },
}


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, color, no_line=True):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if no_line:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, color=None, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    if color:
        p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def build_cover_slide(prs, style, title, subtitle):
    """封面页"""
    s = STYLES[style]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, s["bg"])

    # 顶部装饰条
    add_shape_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), s["accent"])

    # 左侧装饰线
    add_shape_rect(slide, Inches(1.5), Inches(2.8), Inches(2.5), Pt(3), s["accent"])

    # 标题
    add_text_box(slide, title, 1.5, 3.1, 10, 1.4,
                 font_size=44, color=s["text"], bold=True)

    # 副标题
    add_text_box(slide, subtitle, 1.5, 4.6, 10, 0.6,
                 font_size=20, color=s["subtext"])

    # 底部装饰条
    add_shape_rect(slide, Inches(0), Inches(7.2), Inches(13.33), Inches(0.3), s["accent2"])


def build_content_slide(prs, style, title, bullets, slide_num=None, total_slides=None):
    """内容页"""
    s = STYLES[style]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, s["bg"])

    # 左侧装饰条
    add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), s["accent"])

    # 顶部细线
    add_shape_rect(slide, Inches(0.15), Inches(0), Inches(13.18), Inches(0.03), s["accent"])

    # 标题
    add_text_box(slide, title, 0.8, 0.4, 11, 0.9,
                 font_size=32, color=s["accent"], bold=True)

    # 分隔线
    add_shape_rect(slide, Inches(0.8), Inches(1.35), Inches(11.5), Pt(1.5), s["divider"])

    # 要点内容
    y = 1.7
    max_bullets = min(len(bullets), 8)  # 最多8个要点
    bullet_spacing = min(0.7, (5.0 / max(max_bullets, 1)))

    for i in range(max_bullets):
        bullet = bullets[i]

        # 圆点标记
        dot = slide.shapes.add_shape(
            9,  # oval
            Inches(1.0), Inches(y + 0.12),
            Inches(0.13), Inches(0.13)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = s["accent"]
        dot.line.fill.background()

        # 文本 - 根据内容长度调整字号
        font_size = 18 if len(bullet) < 60 else 16
        add_text_box(slide, bullet, 1.35, y, 10.5, 0.6,
                     font_size=font_size, color=s["text"])
        y += bullet_spacing

    # 页码
    if slide_num and total_slides:
        add_text_box(slide, f"{slide_num} / {total_slides}",
                     11.5, 7.0, 1.5, 0.4,
                     font_size=11, color=s["subtext"],
                     alignment=PP_ALIGN.RIGHT)

    return slide


def build_section_slide(prs, style, title, subtitle=""):
    """章节分隔页"""
    s = STYLES[style]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, s["bg"])

    # 中央装饰线
    add_shape_rect(slide, Inches(5.5), Inches(2.8), Inches(2.3), Pt(3), s["accent"])

    # 标题
    add_text_box(slide, title, 0, 3.1, 13.33, 1.2,
                 font_size=40, color=s["text"], bold=True,
                 alignment=PP_ALIGN.CENTER)

    if subtitle:
        add_text_box(slide, subtitle, 0, 4.3, 13.33, 0.6,
                     font_size=18, color=s["subtext"],
                     alignment=PP_ALIGN.CENTER)


def build_end_slide(prs, style):
    """结束页"""
    s = STYLES[style]
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, s["bg"])

    # 装饰线
    add_shape_rect(slide, Inches(5.5), Inches(2.5), Inches(2.3), Pt(3), s["accent"])

    add_text_box(slide, "谢谢", 0, 2.8, 13.33, 1.5,
                 font_size=48, color=s["accent"], bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, "THANK YOU", 0, 4.3, 13.33, 0.8,
                 font_size=20, color=s["subtext"],
                 alignment=PP_ALIGN.CENTER)

    # 底部装饰条
    add_shape_rect(slide, Inches(0), Inches(7.2), Inches(13.33), Inches(0.3), s["accent2"])


def build_ppt(config):
    """根据配置生成 PPT"""
    title = config.get("title", "演示文稿")
    subtitle = config.get("subtitle", "")
    style = config.get("style", "business")
    slides_data = config.get("slides", [])
    output_path = config.get("output", "output.pptx")

    if style not in STYLES:
        style = "business"

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    total_content_slides = len(slides_data)

    # 1. 封面
    build_cover_slide(prs, style, title, subtitle)

    # 2. 目录页（如果有3个以上章节）
    if total_content_slides >= 3:
        toc_bullets = [s.get("title", f"第{i+1}部分") for i, s in enumerate(slides_data)]
        build_content_slide(prs, style, "目录 | CONTENTS", toc_bullets)

    # 3. 内容页
    for i, slide_data in enumerate(slides_data):
        slide_title = slide_data.get("title", f"第{i+1}页")
        bullets = slide_data.get("bullets", [])
        notes = slide_data.get("notes", "")

        # 如果只有一个要点且非常长，拆分为多行
        if len(bullets) == 1 and len(bullets[0]) > 100:
            parts = bullets[0].split("。")
            bullets = [p.strip() + "。" for p in parts if p.strip()]

        slide = build_content_slide(
            prs, style, slide_title, bullets,
            slide_num=i + 1, total_slides=total_content_slides
        )

        # 添加演讲者备注
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    # 4. 结束页
    build_end_slide(prs, style)

    # 确保输出目录存在
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)

    prs.save(output_path)
    print(f"PPT 已保存: {output_path}")
    print(f"共 {len(prs.slides)} 页 (封面 + {'目录 + ' if total_content_slides >= 3 else ''}{total_content_slides} 内容页 + 结束页)")
    print(f"样式: {style}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("用法: python cks_ppt_builder.py <config.json>")
        sys.exit(1)

    config_path = sys.argv[1]
    try:
        with open(config_path, "r", encoding="utf-8") as f:
            config = json.load(f)
    except Exception as e:
        print(f"ERROR: 读取配置文件失败: {e}")
        sys.exit(1)

    build_ppt(config)
